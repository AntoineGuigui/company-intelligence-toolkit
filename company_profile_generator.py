#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
GÉNÉRATEUR DE FICHES COMPAGNIES DÉFENSE - VERSION FINALE v3
✅ Toutes les corrections:
   - [Company Name] en haut du PPT
   - Confidence Index avec étoiles rouges FONCTIONNEL
   - Notes remplies
   - Drapeau du pays en haut à gauche
   - Fiche ouvre automatiquement après génération
"""

import os
import sys
import subprocess
from pathlib import Path
from tkinter import Tk, filedialog, messagebox, simpledialog
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import json
from datetime import datetime

class CompanyProfileGenerator:
    """Générateur de fiches compagnies avec REMPLISSAGE COMPLET"""

    def __init__(self, project_folder):
        """Initialise le générateur"""
        self.project_folder = Path(project_folder)
        self.excel_path = self.project_folder / "DataBase.xlsm"
        self.template_path = self.project_folder / "Template.pptx"
        self.logo_folder = self.project_folder / "logo"
        self.flags_folder = self.project_folder / "Flags"
        self.output_folder = self.project_folder / "Sorties"

        self.output_folder.mkdir(exist_ok=True)
        self._validate_files()

        print("\n🔄 Chargement de l'Excel...")
        self.db = self._load_excel()
        print("✅ Excel chargé avec succès!")
        self.template = None

    def _validate_files(self):
        """Vérifie que tous les fichiers nécessaires existent"""
        if not self.excel_path.exists():
            raise FileNotFoundError(f"Excel non trouvé: {self.excel_path}")
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template non trouvé: {self.template_path}")
        if not self.logo_folder.exists():
            raise FileNotFoundError(f"Dossier logo non trouvé: {self.logo_folder}")
        print(f"✓ Fichiers validés")

    def _load_excel(self):
        """Charge toutes les feuilles Excel"""
        try:
            excel_file = pd.ExcelFile(self.excel_path)

            db = {}
            for sheet in excel_file.sheet_names:
                try:
                    db[sheet] = pd.read_excel(self.excel_path, sheet_name=sheet, dtype=str)
                except Exception as e:
                    db[sheet] = pd.read_excel(self.excel_path, sheet_name=sheet)

            return db

        except Exception as e:
            print(f"❌ Erreur critique lors du chargement: {e}")
            import traceback
            traceback.print_exc()
            raise

    def _find_company(self, company_name, country=None):
        """Trouve une compagnie dans la base de données"""
        df = self.db['DataBase']
        mask = df['Company Name'].str.lower() == company_name.lower()

        if country:
            mask &= df['Country'].str.lower() == country.lower()

        results = df[mask]

        if len(results) == 0:
            print(f"❌ Compagnie '{company_name}' non trouvée")
            return None
        elif len(results) > 1 and not country:
            print(f"⚠ Plusieurs résultats pour '{company_name}'")
            return results.iloc[0].to_dict()
        else:
            return results.iloc[0].to_dict()

    def _get_last_n_years_data(self, sheet_name, company_name, n=4):
        """Récupère les n dernières années de données"""
        try:
            if sheet_name not in self.db:
                return {}

            df = self.db[sheet_name]

            if 'Company' not in df.columns:
                return {}

            mask = df['Company'].str.lower() == company_name.lower()
            row = df[mask]

            if len(row) == 0:
                return {}

            row = row.iloc[0]

            # Extraire toutes les années
            years_data = {}
            for col in df.columns:
                if col != 'Company':
                    try:
                        year = int(col)
                    except (ValueError, TypeError):
                        continue

                    try:
                        value = row[col]

                        if pd.isna(value):
                            continue

                        value_str = str(value).strip()
                        if value_str == '' or value_str.lower() == 'nan':
                            continue

                        try:
                            float_value = float(value_str)
                            if not pd.isna(float_value) and float_value == float_value:
                                years_data[year] = float_value
                        except (ValueError, TypeError):
                            continue

                    except Exception as e:
                        continue

            if not years_data:
                return {}

            sorted_years = sorted(years_data.items(), key=lambda x: x[0])
            sorted_years = sorted_years[-n:]

            return {year: value for year, value in sorted_years}

        except Exception as e:
            print(f"⚠ Erreur dans _get_last_n_years_data: {e}")
            return {}

    def _find_logo(self, company_name):
        """Trouve le logo de la compagnie"""
        clean_name = company_name.lower().replace(" ", "_").replace(".", "")

        for ext in ['.png', '.jpg', '.jpeg', '.gif']:
            logo_path = self.logo_folder / f"{clean_name}{ext}"
            if logo_path.exists():
                return str(logo_path)

            logo_path = self.logo_folder / f"{company_name}{ext}"
            if logo_path.exists():
                return str(logo_path)

        return None

    def _find_flag(self, country_name):
        """Trouve le drapeau du pays"""
        if not self.flags_folder.exists():
            return None

        clean_name = country_name.lower().replace(" ", "_").replace(".", "")

        for ext in ['.png', '.jpg', '.jpeg', '.gif']:
            flag_path = self.flags_folder / f"{clean_name}{ext}"
            if flag_path.exists():
                return str(flag_path)

        # Essayer aussi avec le nom exact
        for ext in ['.png', '.jpg', '.jpeg', '.gif']:
            flag_path = self.flags_folder / f"{country_name}{ext}"
            if flag_path.exists():
                return str(flag_path)

        return None

    def _auto_adjust_text_frame_height(self, shape, text_length):
        """
        Ajuste automatiquement la hauteur du shape selon le contenu

        ✅ Empêche les chevauchements avec les shapes suivants
        ✅ Calcule basé sur la longueur du texte
        ✅ Réduit les marges pour maximiser l'espace
        """
        try:
            if not hasattr(shape, "height"):
                return

            # Estimation: ~65 caractères par ligne à taille standard
            lines_needed = max(1, (text_length // 65) + 1)

            # ✅ RÉDUIT: Hauteur plus conservatrice
            # 0.3 pouces par ligne (au lieu de 0.4)
            new_height = Inches(0.25 + (lines_needed * 0.22))

            # Limiter à une hauteur maximale raisonnable
            max_height = Inches(2.0)  # ← Réduit de 2.5 à 2.0
            if new_height > max_height:
                new_height = max_height

            shape.height = new_height

            # ✅ RÉDUIRE LES MARGES INTERNES du text_frame
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                text_frame.margin_bottom = Inches(0.05)  # Marge basse minimal
                text_frame.margin_top = Inches(0.05)     # Marge haute minimal
                text_frame.margin_left = Inches(0.05)    # Marge gauche minimal
                text_frame.margin_right = Inches(0.05)   # Marge droite minimal
                text_frame.word_wrap = True

                # ✅ RÉDUIRE L'ESPACEMENT ENTRE LES LIGNES
                for paragraph in text_frame.paragraphs:
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(2)  # Minimal
                    paragraph.line_spacing = 1.0   # Simple interlignage

            return new_height
        except Exception as e:
            print(f"  ⚠ Erreur ajustement hauteur: {e}")
            return None

    def _adjust_font_size_for_content(self, text_length, base_size_pt):
        """
        Ajuste la taille de la police selon la longueur du texte
        ✅ Plus agressif pour les textes longs
        """
        # Convertir en points si nécessaire
        if isinstance(base_size_pt, float):
            base_pt = base_size_pt
        else:
            try:
                base_pt = base_size_pt.pt if hasattr(base_size_pt, 'pt') else float(base_size_pt)
            except:
                base_pt = 11

        # Vérifier que base_pt est valide
        if base_pt is None or base_pt <= 0:
            base_pt = 11

        # ✅ AJUSTEMENT PLUS AGRESSIF
        if text_length < 80:
            adjusted = base_pt
        elif text_length < 150:
            adjusted = max(7, base_pt - 1)
        elif text_length < 250:
            adjusted = max(7, base_pt - 2)
        elif text_length < 400:
            adjusted = max(7, base_pt - 3)
        elif text_length < 600:
            adjusted = max(6, base_pt - 4)
        else:
            adjusted = max(6, base_pt - 5)

        return Pt(int(adjusted))

    def _replace_text_with_format(self, shape, old_text, new_text, auto_fit=True):
        """Remplace le texte en préservant la mise en forme avec auto-fit optionnel"""
        if not hasattr(shape, "text_frame"):
            return False

        text_frame = shape.text_frame

        if old_text not in shape.text:
            return False

        if not text_frame.paragraphs:
            return False

        first_para = text_frame.paragraphs[0]
        original_alignment = first_para.alignment

        original_font_props = {}
        original_size_pt = None
        if first_para.runs:
            orig_run = first_para.runs[0]
            try:
                color_rgb = orig_run.font.color.rgb
            except:
                color_rgb = None

            if orig_run.font.size:
                original_size_pt = orig_run.font.size.pt

            original_font_props = {
                'name': orig_run.font.name,
                'size': orig_run.font.size,
                'bold': orig_run.font.bold,
                'italic': orig_run.font.italic,
                'color': color_rgb,
            }

        # Déterminer la taille à utiliser
        if auto_fit and original_size_pt:
            adjusted_size = self._adjust_font_size_for_content(len(new_text), original_size_pt)
        else:
            adjusted_size = None
            if original_font_props.get('size'):
                try:
                    adjusted_size = Pt(int(original_font_props['size'].pt) if hasattr(original_font_props['size'], 'pt') else int(original_font_props['size']))
                except:
                    adjusted_size = Pt(11)

        # Remplacer en parcourant chaque run
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    if adjusted_size:
                        run.font.size = adjusted_size
                    return True

        # Si pas trouvé, reconstruire
        text_frame.clear()
        new_para = text_frame.paragraphs[0]
        new_para.alignment = original_alignment

        run = new_para.add_run()
        run.text = shape.text.replace(old_text, new_text)

        if original_font_props.get('name'):
            run.font.name = original_font_props['name']
        if adjusted_size:
            run.font.size = adjusted_size
        elif original_font_props.get('size'):
            run.font.size = original_font_props['size']
        if original_font_props.get('bold') is not None:
            run.font.bold = original_font_props['bold']
        if original_font_props.get('italic') is not None:
            run.font.italic = original_font_props['italic']
        if original_font_props.get('color'):
            try:
                run.font.color.rgb = original_font_props['color']
            except:
                pass

        try:
            text_frame.word_wrap = True
        except:
            pass

        return True

    def _set_text_with_formatting(self, text_frame, new_text, auto_fit=True):
        """Remplace le texte en préservant la mise en forme avec auto-fit"""
        if not text_frame.paragraphs:
            return

        first_para = text_frame.paragraphs[0]
        original_alignment = first_para.alignment
        original_level = first_para.level

        original_font = None
        original_size = None
        if first_para.runs:
            original_font = first_para.runs[0].font
            if original_font.size:
                original_size = original_font.size.pt

        # Déterminer la taille de police à utiliser
        if auto_fit and original_size:
            adjusted_size = self._adjust_font_size_for_content(len(new_text), original_size)
        else:
            adjusted_size = None
            if original_size:
                try:
                    adjusted_size = Pt(int(original_size))
                except:
                    adjusted_size = Pt(11)

        text_frame.clear()
        new_para = text_frame.paragraphs[0]
        new_para.alignment = original_alignment
        new_para.level = original_level

        run = new_para.add_run()
        run.text = new_text

        if original_font:
            if original_font.name:
                run.font.name = original_font.name
            if adjusted_size:
                run.font.size = adjusted_size
            if original_font.bold is not None:
                run.font.bold = original_font.bold
            if original_font.italic is not None:
                run.font.italic = original_font.italic

            try:
                if original_font.color.rgb:
                    run.font.color.rgb = original_font.color.rgb
            except:
                pass

        try:
            text_frame.word_wrap = True
        except:
            pass

    def _fill_template(self, company_data, company_name):
        """Remplit le template PowerPoint COMPLÈTEMENT"""
        prs = Presentation(str(self.template_path))
        slide = prs.slides[0]

        print("\n📝 Remplissage du template...")

        # ============================================
        # 1. TITRE ET INFORMATIONS BASIQUES
        # ============================================
        print("✓ Remplissage du titre...")

        # ✅ REMPLIR LE TITRE [Company Name] DIRECTEMENT (c'est un PLACEHOLDER special)
        try:
            company_name_text = str(company_data.get('Company Name', ''))
            for shape in slide.shapes:
                if hasattr(shape, "name") and shape.name == "Titre 1":
                    if hasattr(shape, "text_frame"):
                        text_frame = shape.text_frame
                        if text_frame.paragraphs:
                            para = text_frame.paragraphs[0]
                            para.text = company_name_text
                            print(f"  ✓ Titre remplacé: {company_name_text}")
                        break
        except Exception as e:
            print(f"  ⚠ Erreur titre: {e}")

        # Remplir les autres placeholders
        replacements = {
            '[Country]': company_data.get('Country', ''),
            '[Activity]': company_data.get('Activity', ''),
            '[Business Overview]': company_data.get('Business Overview', ''),
            '[Business relationships]': company_data.get('Business relationships', ''),
            '[Capability]': company_data.get('Capability', ''),
        }

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for placeholder, value in replacements.items():
                    if placeholder in shape.text:
                        # ✅ Remplacer le texte
                        self._replace_text_with_format(shape, placeholder, str(value))
                        # ✅ AJUSTER LA HAUTEUR pour éviter les chevauchements
                        self._auto_adjust_text_frame_height(shape, len(str(value)))
                        print(f"  → Remplacé + hauteur ajustée: {placeholder}")

        # ============================================
        # 2. AJOUTER LE DRAPEAU EN HAUT À GAUCHE
        # ============================================
        print("✓ Ajout du drapeau et alignement...")

        country = company_data.get('Country', '')
        flag_path = self._find_flag(country)

        # Trouver la position du shape "[Country] Defence Companies"
        country_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text") and "[Country] Defence Companies" in shape.text:
                country_shape = shape
                break

        if flag_path:
            try:
                # ✅ Position du drapeau
                flag_width = Inches(0.35)
                flag_height = Inches(0.25)

                # ✅ Positionner le drapeau à gauche
                # Aligner avec le texte "[Country] Defence Companies"
                if country_shape:
                    # Drapeau juste à gauche du texte, au même niveau vertical
                    flag_left = Inches(0.15)
                    flag_top = country_shape.top + (country_shape.height - flag_height) / 2  # Centre verticalement
                else:
                    # Position par défaut si le shape n'est pas trouvé
                    flag_left = Inches(0.15)
                    flag_top = Inches(0.15)

                slide.shapes.add_picture(flag_path, flag_left, flag_top, flag_width, flag_height)
                print(f"  → Drapeau aligné avec texte: {country}")

                # ✅ Ajuster la position du texte "[Country] Defence Companies"
                # pour qu'il soit à droite du drapeau
                if country_shape:
                    # Laisser de l'espace pour le drapeau
                    new_left = flag_left + flag_width + Inches(0.1)  # 0.1" d'espacement

                    # Si le texte chevauche, le déplacer
                    if country_shape.left < new_left + Inches(0.5):
                        country_shape.left = new_left
                        print(f"  → Texte déplacé à droite du drapeau")

            except Exception as e:
                print(f"  ⚠ Erreur ajout/alignement drapeau: {e}")
                import traceback
                traceback.print_exc()

        # ============================================
        # 3. TABLEAU "KEY FACTS" (Tableau 1)
        # ============================================
        print("✓ Remplissage du tableau Key facts...")

        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
                table = shape.table

                if table_count == 1:
                    if len(table.rows) >= 3 and len(table.columns) >= 2:
                        cell = table.cell(0, 1)
                        if cell.text_frame:
                            self._set_text_with_formatting(cell.text_frame, str(company_data.get('Field', '')))

                        cell = table.cell(1, 1)
                        if cell.text_frame:
                            self._set_text_with_formatting(cell.text_frame, str(company_data.get('Locations', '')))

                        cell = table.cell(2, 1)
                        if cell.text_frame:
                            self._set_text_with_formatting(cell.text_frame, str(company_data.get('Founded', '')))

                        if len(table.rows) >= 4:
                            cell = table.cell(3, 1)
                            if cell.text_frame:
                                self._set_text_with_formatting(cell.text_frame, str(company_data.get('N° employees', '')))
                        if len(table.rows) >= 5:
                            cell = table.cell(4, 1)
                            if cell.text_frame:
                                self._set_text_with_formatting(cell.text_frame, str(company_data.get('Key people', '')))
                        if len(table.rows) >= 6:
                            cell = table.cell(5, 1)
                            if cell.text_frame:
                                self._set_text_with_formatting(cell.text_frame, str(company_data.get('Type Ownership', '')))

        # ============================================
        # 4. TABLEAU FINANCIERS (Tableau 2)
        # ============================================
        print("✓ Remplissage du tableau Financials...")

        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
                table = shape.table

                if table_count == 2:
                    revenue_data = self._get_last_n_years_data('Revenue', company_name, n=4)
                    ebit_data = self._get_last_n_years_data('EBIT', company_name, n=4)
                    net_profit_data = self._get_last_n_years_data('Net Profit', company_name, n=4)
                    ebit_margin_data = self._get_last_n_years_data('EBIT Margin', company_name, n=4)
                    net_profit_margin_data = self._get_last_n_years_data('Net Profit Margin', company_name, n=4)

                    years_list = sorted(revenue_data.keys()) if revenue_data else []

                    # En-têtes
                    for col_idx, year in enumerate(years_list, start=1):
                        if col_idx < len(table.columns):
                            cell = table.cell(0, col_idx)
                            if cell.text_frame:
                                self._set_text_with_formatting(cell.text_frame, f"FY{year % 100}")

                    if len(table.rows) >= 3 and len(revenue_data) > 0:
                        # Ligne 2: Revenues
                        for col_idx, year in enumerate(years_list, start=1):
                            if col_idx < len(table.columns) and year in revenue_data:
                                val = revenue_data[year]
                                formatted = f"{int(val) if val == int(val) else f'{val:.1f}'}"
                                cell = table.cell(2, col_idx)
                                if cell.text_frame:
                                    self._set_text_with_formatting(cell.text_frame, formatted)

                        # Ligne 4: EBIT
                        if len(table.rows) >= 5:
                            for col_idx, year in enumerate(years_list, start=1):
                                if col_idx < len(table.columns) and year in ebit_data:
                                    val = ebit_data[year]
                                    formatted = f"{int(val) if val == int(val) else f'{val:.1f}'}"
                                    cell = table.cell(4, col_idx)
                                    if cell.text_frame:
                                        self._set_text_with_formatting(cell.text_frame, formatted)

                        # Ligne 5: EBIT Margin
                        if len(table.rows) >= 6:
                            for col_idx, year in enumerate(years_list, start=1):
                                if col_idx < len(table.columns) and year in ebit_margin_data:
                                    val = ebit_margin_data[year]
                                    formatted = f"{int(val) if val == int(val) else f'{val:.1f}'}%"
                                    cell = table.cell(5, col_idx)
                                    if cell.text_frame:
                                        self._set_text_with_formatting(cell.text_frame, formatted)

                        # Ligne 7: Net Profit
                        if len(table.rows) >= 8:
                            for col_idx, year in enumerate(years_list, start=1):
                                if col_idx < len(table.columns) and year in net_profit_data:
                                    val = net_profit_data[year]
                                    formatted = f"{int(val) if val == int(val) else f'{val:.1f}'}"
                                    cell = table.cell(7, col_idx)
                                    if cell.text_frame:
                                        self._set_text_with_formatting(cell.text_frame, formatted)

                        # Ligne 8: Net Profit Margin
                        if len(table.rows) >= 9:
                            for col_idx, year in enumerate(years_list, start=1):
                                if col_idx < len(table.columns) and year in net_profit_margin_data:
                                    val = net_profit_margin_data[year]
                                    formatted = f"{int(val) if val == int(val) else f'{val:.1f}'}%"
                                    cell = table.cell(8, col_idx)
                                    if cell.text_frame:
                                        self._set_text_with_formatting(cell.text_frame, formatted)

        # ============================================
        # 5. CONFIDENCE INDEX (ÉTOILES ROUGES)
        # ============================================
        print("✓ Remplissage du Confidence Index...")

        try:
            confidence_str = str(company_data.get('Confidence Index', '0')).strip()
            confidence = int(confidence_str) if confidence_str.isdigit() else 0
            confidence = max(1, min(5, confidence))

            print(f"  → Confidence Index: {confidence}/5")

            # Les étoiles sont dans un groupe imbriqué: Group 27 > Group 30 > 5 étoiles
            for shape in slide.shapes:
                if shape.name == "Groupe 27":
                    # C'est le groupe principal des étoiles
                    if hasattr(shape, 'shapes'):
                        for sub_shape in shape.shapes:
                            if sub_shape.name == "Groupe 30":
                                # C'est le sous-groupe avec les 5 étoiles
                                if hasattr(sub_shape, 'shapes'):
                                    for star_idx, star_shape in enumerate(sub_shape.shapes):
                                        if "Étoile" in star_shape.name:
                                            # Remplir les N premières étoiles en rouge
                                            if star_idx < confidence:
                                                star_shape.fill.solid()
                                                star_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
                                                print(f"    ⭐ Étoile {star_idx + 1} remplie en rouge")
                                            else:
                                                star_shape.fill.solid()
                                                star_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
                                                print(f"    ⭐ Étoile {star_idx + 1} grise")

        except Exception as e:
            print(f"  ⚠ Erreur Confidence Index: {e}")
            import traceback
            traceback.print_exc()


        # ============================================
        # 6. AJOUTER LES SPEAKER NOTES (notes sous la slide)
        # ============================================
        print("✓ Ajout des Speaker Notes...")

        try:
            notes_text = ""

            # Ajouter les informations pertinentes aux notes

            if company_data.get('Notes'):
                notes_text += f"\nNotes:\n{company_data.get('Notes')}\n"

            # ✅ Ajouter les notes à la slide
            if notes_text.strip():
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes_text.strip()
                print(f"  → Speaker Notes ajoutées")

        except Exception as e:
            print(f"  ⚠ Erreur Speaker Notes: {e}")
            import traceback
            traceback.print_exc()
        # ============================================
        logo_path = self._find_logo(company_name)
        if logo_path:
            try:
                print("✓ Ajout du logo...")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and "Logo" in shape.text:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        sp = shape.element
                        sp.getparent().remove(sp)
                        try:
                            slide.shapes.add_picture(logo_path, left, top, width, height)
                        except:
                            slide.shapes.add_picture(logo_path, left, top)
                        break
            except Exception as e:
                print(f"  ⚠ Erreur logo: {e}")

        return prs

    def generate_profile(self, company_name, country=None):
        """Génère une fiche PowerPoint pour une compagnie"""
        print(f"\n{'='*60}")
        print(f"Génération de la fiche: {company_name}")
        print(f"{'='*60}")

        company_data = self._find_company(company_name, country)
        if company_data is None:
            return None

        print(f"✓ Compagnie trouvée: {company_data.get('Company Name')}")

        prs = self._fill_template(company_data, company_name)

        output_filename = f"{company_name.replace(' ', '_')}.pptx"
        output_path = self.output_folder / output_filename

        # ✅ ÉCRASE LA VERSION PRÉCÉDENTE
        prs.save(str(output_path))
        print(f"✓ Fiche sauvegardée: {output_path}")

        return str(output_path)

    def list_companies(self):
        """Liste toutes les compagnies disponibles"""
        df = self.db['DataBase']
        companies = df[['Country', 'Company Name']].drop_duplicates()
        return companies

def find_project_folder():
    """Cherche le dossier Projet automatiquement"""

    current_path = Path.cwd()

    for parent in [current_path] + list(current_path.parents):
        if (parent / "DataBase.xlsm").exists() and (parent / "Template.pptx").exists():
            print(f"✓ Dossier Projet trouvé: {parent}")
            return str(parent)

    print("⚠ Dossier Projet non trouvé automatiquement")
    root = Tk()
    root.withdraw()
    folder = filedialog.askdirectory(title="Sélectionner le dossier 'Projet'")
    root.destroy()

    if folder:
        return folder

    return None

def open_file(file_path):
    """Ouvre le fichier avec l'application par défaut"""
    try:
        if sys.platform == 'win32':
            os.startfile(file_path)
        elif sys.platform == 'darwin':
            subprocess.Popen(['open', file_path])
        else:
            subprocess.Popen(['xdg-open', file_path])
        print(f"✓ Fichier ouvert: {file_path}")
    except Exception as e:
        print(f"⚠ Erreur ouverture fichier: {e}")

def main():
    """Fonction principale"""
    print("\n" + "="*80)
    print("🚀 GÉNÉRATEUR DE FICHES COMPAGNIES DÉFENSE - VERSION FINALE v3")
    print("="*80)

    project_folder = find_project_folder()
    if not project_folder:
        print("❌ Pas de dossier sélectionné")
        return

    try:
        print(f"\n📂 Dossier Projet: {project_folder}")

        generator = CompanyProfileGenerator(project_folder)

        print("\n" + "="*80)
        print("✅ PRÊT À GÉNÉRER LES FICHES!")
        print("="*80)

        while True:
            root = Tk()
            root.withdraw()

            messagebox.showinfo(
                "Menu",
                "Choisir une option:\n\n"
                "1 - Générer une fiche\n"
                "2 - Lister les compagnies\n"
                "3 - Quitter"
            )

            choice = simpledialog.askstring(
                "Menu",
                "Entrer le numéro (1, 2 ou 3):"
            )
            root.destroy()

            try:
                choice = int(choice) if choice else None
            except ValueError:
                choice = None

            if choice == 1:
                root = Tk()
                root.withdraw()
                company_name = simpledialog.askstring(
                    "Générer une fiche",
                    "Entrer le nom de la compagnie:"
                )
                root.destroy()

                if company_name:
                    result = generator.generate_profile(company_name)

                    if result:
                        # ✅ OUVRIR AUTOMATIQUEMENT LE FICHIER
                        open_file(result)

                        root = Tk()
                        root.withdraw()
                        messagebox.showinfo(
                            "✅ Succès",
                            f"Fiche générée et ouverture en cours..."
                        )
                        root.destroy()

            elif choice == 2:
                companies = generator.list_companies()
                root = Tk()
                root.withdraw()
                messagebox.showinfo(
                    "Compagnies disponibles",
                    f"Total: {len(companies)} compagnies"
                )
                root.destroy()

            elif choice == 3 or choice is None:
                print("\n✅ Fermeture")
                break

    except Exception as e:
        print(f"❌ Erreur: {e}")
        import traceback
        traceback.print_exc()

        root = Tk()
        root.withdraw()
        messagebox.showerror("❌ Erreur", f"Une erreur s'est produite:\n{e}")
        root.destroy()

if __name__ == "__main__":
    main()
